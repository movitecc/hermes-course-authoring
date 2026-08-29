#!/usr/bin/env python3
"""Local, dependency-light course lifecycle and export engine for Hermes."""
from __future__ import annotations

import argparse
import datetime as dt
import html
import json
import os
import re
import shutil
import sys
import tempfile
import uuid
import zipfile
from pathlib import Path
from typing import Any

SCHEMA_VERSION = 1


def now() -> str:
    return dt.datetime.now(dt.timezone.utc).replace(microsecond=0).isoformat()


def slug(value: str) -> str:
    value = re.sub(r"[^a-zA-Z0-9\u0080-\uffff]+", "-", value.strip().lower()).strip("-")
    return value or "course"


def atomic_json(path: Path, value: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    fd, name = tempfile.mkstemp(prefix=f".{path.name}.", dir=path.parent)
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            json.dump(value, f, ensure_ascii=False, indent=2)
            f.write("\n")
        os.replace(name, path)
    finally:
        if os.path.exists(name):
            os.unlink(name)


def workspace(args: argparse.Namespace) -> Path:
    return Path(args.workspace).expanduser().resolve()


def course_path(root: Path, course_id: str) -> Path:
    return root / "courses" / course_id / "course.json"


def load_course(root: Path, course_id: str) -> dict[str, Any]:
    path = course_path(root, course_id)
    if not path.exists():
        raise SystemExit(f"course not found: {course_id}")
    return json.loads(path.read_text(encoding="utf-8"))


def save_course(root: Path, course: dict[str, Any]) -> None:
    course["updatedAt"] = now()
    atomic_json(course_path(root, course["id"]), course)


def read_payload(path: str) -> Any:
    p = Path(path).expanduser()
    if not p.exists():
        raise SystemExit(f"payload not found: {p}")
    if p.suffix.lower() == ".json":
        return json.loads(p.read_text(encoding="utf-8"))
    return {"content": p.read_text(encoding="utf-8")}


def cmd_init(args: argparse.Namespace) -> None:
    root = workspace(args)
    course_id = f"{slug(args.title)}-{uuid.uuid4().hex[:8]}"
    course = {
        "schemaVersion": SCHEMA_VERSION,
        "id": course_id,
        "title": args.title,
        "topic": args.topic,
        "audience": args.audience,
        "language": args.language,
        "durationMinutes": args.duration,
        "style": args.style,
        "objectives": [],
        "prerequisites": [],
        "sources": [],
        "lessons": [],
        "sessions": [],
        "progress": {"completedLessons": [], "weakAreas": [], "scores": []},
        "artifacts": [],
        "createdAt": now(),
        "updatedAt": now(),
        "versions": [{"at": now(), "reason": "created", "lessonCount": 0}],
    }
    save_course(root, course)
    print(json.dumps({"courseId": course_id, "path": str(course_path(root, course_id))}, ensure_ascii=False))


def cmd_lesson(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    payload = read_payload(args.content)
    lesson = payload if isinstance(payload, dict) else {"content": str(payload)}
    lesson.setdefault("id", args.lesson)
    lesson.setdefault("title", args.lesson)
    lesson.setdefault("objectives", [])
    lesson.setdefault("keyConcepts", [])
    lesson.setdefault("scenes", [{"id": f"{args.lesson}-scene-1", "title": lesson["title"], "content": lesson.get("content", "")}])
    lesson.setdefault("exercises", [])
    lesson.setdefault("version", 1)
    lesson["updatedAt"] = now()
    existing = next((i for i, x in enumerate(course["lessons"]) if x["id"] == args.lesson), None)
    if existing is None:
        course["lessons"].append(lesson)
    else:
        lesson["version"] = int(course["lessons"][existing].get("version", 1)) + 1
        course["lessons"][existing] = lesson
    course["versions"].append({"at": now(), "reason": f"lesson:{args.lesson}", "lessonCount": len(course["lessons"])})
    save_course(root, course)
    print(json.dumps({"courseId": course["id"], "lessonId": args.lesson, "version": lesson["version"]}, ensure_ascii=False))


def cmd_scene(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    lesson = next((x for x in course["lessons"] if x["id"] == args.lesson), None)
    if lesson is None:
        raise SystemExit(f"lesson not found: {args.lesson}")
    payload = read_payload(args.scene)
    scene = payload.get("scene", payload) if isinstance(payload, dict) else payload
    if not isinstance(scene, dict) or not scene.get("id"):
        raise SystemExit("scene payload must be an object with id")
    lesson.setdefault("scenes", [])
    old = next((i for i, x in enumerate(lesson["scenes"]) if x.get("id") == scene["id"]), None)
    if old is None:
        lesson["scenes"].append(scene)
    else:
        lesson["scenes"][old] = scene
    lesson["version"] = int(lesson.get("version", 1)) + 1
    save_course(root, course)
    print(json.dumps({"courseId": course["id"], "lessonId": args.lesson, "sceneId": scene["id"], "version": lesson["version"]}, ensure_ascii=False))


def cmd_exercise(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    payload = read_payload(args.exercise)
    exercise = payload if isinstance(payload, dict) else {"prompt": str(payload)}
    exercise.setdefault("id", args.exercise_id)
    exercise.setdefault("type", "short-answer")
    exercise.setdefault("prompt", "Answer the question.")
    exercise.setdefault("answer", "")
    exercise.setdefault("rubric", [])
    lesson = next((x for x in course["lessons"] if x["id"] == args.lesson), None)
    if lesson is None:
        raise SystemExit(f"lesson not found: {args.lesson}")
    lesson.setdefault("exercises", [])
    old = next((i for i, x in enumerate(lesson["exercises"]) if x["id"] == args.exercise_id), None)
    if old is None:
        lesson["exercises"].append(exercise)
    else:
        lesson["exercises"][old] = exercise
    save_course(root, course)
    print(json.dumps({"courseId": course["id"], "lessonId": args.lesson, "exerciseId": exercise["id"]}, ensure_ascii=False))


def cmd_session(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    session_id = f"session-{uuid.uuid4().hex[:10]}"
    session = {"id": session_id, "state": "teaching", "lessonId": args.lesson or (course["lessons"][0]["id"] if course["lessons"] else None), "turn": 0, "answers": [], "startedAt": now(), "updatedAt": now()}
    course["sessions"].append(session)
    save_course(root, course)
    print(json.dumps(session, ensure_ascii=False))


def cmd_list(args: argparse.Namespace) -> None:
    root = workspace(args)
    base = root / "courses"
    rows = []
    if base.exists():
        for path in sorted(base.glob("*/course.json")):
            c = json.loads(path.read_text(encoding="utf-8"))
            rows.append({"id": c["id"], "title": c["title"], "lessons": len(c.get("lessons", [])), "updatedAt": c.get("updatedAt")})
    print(json.dumps(rows, ensure_ascii=False))


def cmd_pause(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    session = next((x for x in course["sessions"] if x["id"] == args.session), None)
    if session is None:
        raise SystemExit(f"session not found: {args.session}")
    session["state"] = "paused"
    session["updatedAt"] = now()
    save_course(root, course)
    print(json.dumps({"sessionId": args.session, "state": "paused"}, ensure_ascii=False))


def cmd_resume(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    session = next((x for x in course["sessions"] if x["id"] == args.session), None)
    if session is None:
        raise SystemExit(f"session not found: {args.session}")
    session["state"] = "teaching"
    session["updatedAt"] = now()
    save_course(root, course)
    print(json.dumps(session, ensure_ascii=False))


def cmd_progress(args: argparse.Namespace) -> None:
    course = load_course(workspace(args), args.course)
    print(json.dumps(course["progress"], ensure_ascii=False))


def cmd_source(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    source = {"url": args.url, "title": args.title or args.url, "addedAt": now()}
    if not any(x.get("url") == args.url for x in course["sources"]):
        course["sources"].append(source)
    save_course(root, course)
    print(json.dumps(source, ensure_ascii=False))


def cmd_artifact(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    path = Path(args.path).expanduser().resolve()
    if not path.exists() or not path.is_file():
        raise SystemExit(f"artifact not found: {path}")
    artifact = {"kind": args.kind, "path": str(path), "bytes": path.stat().st_size, "createdAt": now()}
    course["artifacts"].append(artifact)
    save_course(root, course)
    print(json.dumps(artifact, ensure_ascii=False))


def cmd_answer(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    session = next((x for x in course["sessions"] if x["id"] == args.session), None)
    if session is None:
        raise SystemExit(f"session not found: {args.session}")
    lesson = next((x for x in course["lessons"] if x["id"] == session["lessonId"]), None)
    exercise = next((x for x in (lesson or {}).get("exercises", []) if x["id"] == args.exercise), None)
    if exercise is None:
        raise SystemExit(f"exercise not found: {args.exercise}")
    expected = str(exercise.get("answer", "")).strip().lower()
    actual = args.answer.strip().lower()
    correct = bool(expected) and (actual == expected or expected in actual)
    feedback = "Correct. Nice work." if correct else f"Not quite. Review the lesson and compare with: {exercise.get('answer', 'the rubric').strip()}"
    result = {"exerciseId": args.exercise, "answer": args.answer, "correct": correct, "feedback": feedback, "at": now()}
    session["answers"].append(result)
    session["state"] = "reviewing"
    session["turn"] += 1
    session["updatedAt"] = now()
    if correct and lesson and lesson["id"] not in course["progress"]["completedLessons"]:
        course["progress"]["completedLessons"].append(lesson["id"])
    if not correct:
        for area in lesson.get("keyConcepts", []) if lesson else []:
            if area not in course["progress"]["weakAreas"]:
                course["progress"]["weakAreas"].append(area)
    course["progress"]["scores"].append({"lessonId": session["lessonId"], "exerciseId": args.exercise, "correct": correct, "at": now()})
    save_course(root, course)
    print(json.dumps(result, ensure_ascii=False))


def cmd_pbl_init(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    project = {
        "id": args.project, "title": args.title, "description": args.description, "status": "active",
        "proficiency": {"level": "beginner", "score": 0, "confidence": 0},
        "milestones": [{"id": f"{args.project}-m1", "title": "Define the problem", "status": "active", "submissions": []}],
        "submissions": [], "evaluations": [], "events": [{"kind": "project_started", "at": now()}],
    }
    course.setdefault("pblProjects", []).append(project)
    save_course(root, course)
    print(json.dumps(project, ensure_ascii=False))


def cmd_pbl_submit(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    project = next((x for x in course.get("pblProjects", []) if x["id"] == args.project), None)
    if project is None: raise SystemExit(f"PBL project not found: {args.project}")
    milestone = next((x for x in project["milestones"] if x["id"] == args.milestone), None)
    if milestone is None: raise SystemExit(f"PBL milestone not found: {args.milestone}")
    submission = {"id": f"submission-{uuid.uuid4().hex[:10]}", "milestoneId": args.milestone, "kind": args.kind, "content": args.content, "createdAt": now()}
    milestone.setdefault("submissions", []).append(submission); project["submissions"].append(submission)
    project["events"].append({"kind": "submission_added", "milestoneId": args.milestone, "at": submission["createdAt"]})
    save_course(root, course)
    print(json.dumps(submission, ensure_ascii=False))


def cmd_pbl_complete(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    project = next((x for x in course.get("pblProjects", []) if x["id"] == args.project), None)
    if project is None: raise SystemExit(f"PBL project not found: {args.project}")
    milestone = next((x for x in project["milestones"] if x["id"] == args.milestone), None)
    if milestone is None: raise SystemExit(f"PBL milestone not found: {args.milestone}")
    if not milestone.get("submissions"): raise SystemExit("milestone requires at least one submission")
    milestone["status"] = "completed"
    project["proficiency"]["score"] = min(100, int(project["proficiency"].get("score", 0)) + 25)
    project["proficiency"]["confidence"] = min(100, int(project["proficiency"].get("confidence", 0)) + 20)
    if project["proficiency"]["score"] >= 50: project["proficiency"]["level"] = "intermediate"
    project["events"].append({"kind": "milestone_completed", "milestoneId": args.milestone, "at": now()})
    if all(x.get("status") == "completed" for x in project["milestones"]): project["status"] = "completed"
    save_course(root, course)
    print(json.dumps({"projectId": args.project, "milestoneId": args.milestone, "milestoneStatus": milestone["status"], "projectStatus": project["status"], "proficiency": project["proficiency"]}, ensure_ascii=False))


def _pbl_project(course: dict[str, Any], project_id: str) -> dict[str, Any]:
    project = next((x for x in course.get("pblProjects", []) if x["id"] == project_id), None)
    if project is None: raise SystemExit(f"PBL project not found: {project_id}")
    return project


def _unlock_next_milestone(project: dict[str, Any], milestone_id: str) -> None:
    for i, milestone in enumerate(project["milestones"]):
        if milestone["id"] == milestone_id and i + 1 < len(project["milestones"]):
            if project["milestones"][i + 1].get("status") == "locked":
                project["milestones"][i + 1]["status"] = "active"
            return


def cmd_pbl_plan(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    project = _pbl_project(course, args.project)
    payload = read_payload(args.milestones)
    items = payload.get("milestones", payload) if isinstance(payload, dict) else payload
    if not isinstance(items, list) or not items: raise SystemExit("milestones must be a non-empty list")
    project["milestones"] = [{"id": x.get("id", f"{args.project}-m{i}"), "title": x.get("title", f"Milestone {i}"), "description": x.get("description", ""), "status": "active" if i == 1 else "locked", "submissions": [], "skills": x.get("skills", [])} for i, x in enumerate(items, 1)]
    project["events"].append({"kind": "plan_created", "at": now(), "milestoneCount": len(project["milestones"])})
    save_course(root, course)
    print(json.dumps({"projectId": args.project, "milestones": project["milestones"]}, ensure_ascii=False))


def cmd_pbl_evaluate(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    project = _pbl_project(course, args.project)
    milestone = next((x for x in project["milestones"] if x["id"] == args.milestone), None)
    if milestone is None: raise SystemExit(f"PBL milestone not found: {args.milestone}")
    score = max(0, min(100, args.score))
    evaluation = {"id": f"evaluation-{uuid.uuid4().hex[:10]}", "kind": "milestone", "milestoneId": args.milestone, "score": score, "feedback": args.feedback, "strengths": args.strength or [], "improvements": args.improvement or [], "createdAt": now()}
    project.setdefault("evaluations", []).append(evaluation)
    milestone["evaluation"] = evaluation
    if score >= 50:
        milestone["status"] = "completed"
        _unlock_next_milestone(project, args.milestone)
    scores = [x["score"] for x in project["evaluations"] if isinstance(x.get("score"), (int, float))]
    project["proficiency"]["score"] = round(sum(scores) / len(scores)) if scores else 0
    project["proficiency"]["confidence"] = min(100, len(scores) * 25)
    project["proficiency"]["level"] = "advanced" if project["proficiency"]["score"] >= 80 else "intermediate" if project["proficiency"]["score"] >= 50 else "beginner"
    project["events"].append({"kind": "milestone_evaluated", "milestoneId": args.milestone, "score": score, "at": evaluation["createdAt"]})
    if all(x.get("status") == "completed" for x in project["milestones"]): project["status"] = "completed"
    save_course(root, course)
    print(json.dumps({"evaluation": evaluation, "proficiency": project["proficiency"]}, ensure_ascii=False))


def cmd_pbl_report(args: argparse.Namespace) -> None:
    course = load_course(workspace(args), args.course)
    project = _pbl_project(course, args.project)
    report = {"projectId": project["id"], "title": project["title"], "status": project["status"], "proficiency": project["proficiency"], "milestones": [{"id": x["id"], "title": x["title"], "status": x["status"], "submissionCount": len(x.get("submissions", [])), "score": (x.get("evaluation") or {}).get("score")} for x in project["milestones"]], "evaluations": len(project.get("evaluations", [])), "events": len(project.get("events", []))}
    print(json.dumps(report, ensure_ascii=False, indent=2))


def cmd_pbl_next(args: argparse.Namespace) -> None:
    course = load_course(workspace(args), args.course)
    project = _pbl_project(course, args.project)
    for milestone in project["milestones"]:
        if milestone.get("status") != "active": continue
        for task in milestone.get("microtasks", []):
            if task.get("status", "active") in {"active", "in_progress", "todo"} and not task.get("submission"):
                print(json.dumps({"projectId": project["id"], "milestoneId": milestone["id"], "milestoneTitle": milestone["title"], "microtask": task}, ensure_ascii=False))
                return
    print(json.dumps({"projectId": project["id"], "status": project["status"], "next": None}, ensure_ascii=False))


def cmd_pbl_task_submit(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    project = _pbl_project(course, args.project)
    milestone = next((x for x in project["milestones"] if x["id"] == args.milestone), None)
    if milestone is None: raise SystemExit(f"PBL milestone not found: {args.milestone}")
    task = next((x for x in milestone.get("microtasks", []) if x.get("id") == args.microtask), None)
    if task is None: raise SystemExit(f"PBL microtask not found: {args.microtask}")
    task["submission"] = {"kind": args.kind, "content": args.content, "createdAt": now()}
    task["status"] = "submitted"
    project["events"].append({"kind": "microtask_submitted", "milestoneId": args.milestone, "microtaskId": args.microtask, "at": now()})
    save_course(root, course)
    print(json.dumps({"projectId": args.project, "milestoneId": args.milestone, "microtaskId": args.microtask, "status": task["status"]}, ensure_ascii=False))


def cmd_pbl_task_complete(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    project = _pbl_project(course, args.project)
    milestone = next((x for x in project["milestones"] if x["id"] == args.milestone), None)
    if milestone is None: raise SystemExit(f"PBL milestone not found: {args.milestone}")
    task = next((x for x in milestone.get("microtasks", []) if x.get("id") == args.microtask), None)
    if task is None: raise SystemExit(f"PBL microtask not found: {args.microtask}")
    if not task.get("submission"): raise SystemExit("microtask requires a submission")
    task["status"] = "completed"
    project["events"].append({"kind": "microtask_completed", "milestoneId": args.milestone, "microtaskId": args.microtask, "at": now()})
    if all(x.get("status") == "completed" for x in milestone.get("microtasks", [])):
        milestone["status"] = "completed"
        _unlock_next_milestone(project, args.milestone)
    if all(x.get("status") == "completed" for x in project["milestones"]): project["status"] = "completed"
    save_course(root, course)
    print(json.dumps({"projectId": args.project, "milestoneId": args.milestone, "microtaskId": args.microtask, "milestoneStatus": milestone["status"], "projectStatus": project["status"]}, ensure_ascii=False))


def cmd_pbl_import(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    payload = read_payload(args.project)
    project = payload.get("project", payload) if isinstance(payload, dict) else payload
    if not isinstance(project, dict) or not project.get("milestones"):
        raise SystemExit("planner payload must contain a project with milestones")
    project.setdefault("id", args.project_id)
    project.setdefault("status", "active")
    project.setdefault("submissions", [])
    project.setdefault("evaluations", [])
    project.setdefault("events", []).append({"kind": "planner_imported", "at": now()})
    course.setdefault("pblProjects", [])
    old = next((i for i, x in enumerate(course["pblProjects"]) if x.get("id") == project["id"]), None)
    if old is None: course["pblProjects"].append(project)
    else: course["pblProjects"][old] = project
    save_course(root, course)
    print(json.dumps({"courseId": course["id"], "projectId": project["id"], "milestones": len(project["milestones"]), "microtasks": sum(len(x.get("microtasks", [])) for x in project["milestones"])}, ensure_ascii=False))


def _session(course: dict[str, Any], session_id: str) -> dict[str, Any]:
    session = next((x for x in course["sessions"] if x["id"] == session_id), None)
    if session is None:
        raise SystemExit(f"session not found: {session_id}")
    return session


def cmd_turn(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    session = _session(course, args.session)
    transitions = {
        "lecture": ("teaching", "Explain the current scene, then ask one comprehension question."),
        "question": ("awaiting_answer", "Wait for the learner answer before continuing."),
        "discussion": ("discussing", "Invite teacher, assistant, and peer perspectives; synthesize after responses."),
        "review": ("reviewing", "Review the learner's answer and identify a weak area or next step."),
        "complete": ("completed", "Summarize outcomes and recommend the next lesson."),
    }
    state, next_action = transitions[args.mode]
    session["state"] = state
    session["turn"] = int(session.get("turn", 0)) + 1
    session.setdefault("transcript", []).append({"mode": args.mode, "message": args.message or "", "at": now()})
    session["updatedAt"] = now()
    save_course(root, course)
    print(json.dumps({"sessionId": session["id"], "state": state, "turn": session["turn"], "lessonId": session.get("lessonId"), "nextAction": next_action}, ensure_ascii=False))


def cmd_quiz_answer(args: argparse.Namespace) -> None:
    root = workspace(args)
    course = load_course(root, args.course)
    session = _session(course, args.session)
    scene = next((s for lesson in course["lessons"] for s in lesson.get("scenes", []) if s.get("id") == args.scene), None)
    questions = (((scene or {}).get("content") or {}).get("questions") or [])
    if not questions:
        raise SystemExit(f"quiz scene not found or has no questions: {args.scene}")
    try:
        index = int(args.question)
        question = questions[index]
    except (ValueError, IndexError):
        question = next((q for q in questions if q.get("id") == args.question), None)
    if question is None:
        raise SystemExit(f"quiz question not found: {args.question}")
    expected = question.get("answer", question.get("correctAnswer", ""))
    actual = args.answer
    def norm(v: Any) -> list[str]:
        values = v if isinstance(v, list) else [v]
        return sorted(str(x).strip().lower() for x in values)
    correct = bool(expected) and norm(expected) == norm(actual)
    result = {"sceneId": args.scene, "question": args.question, "answer": actual, "correct": correct, "feedback": "Correct." if correct else "Review the explanation and try again.", "at": now()}
    session.setdefault("quizAnswers", []).append(result)
    session["state"] = "reviewing"
    session["turn"] = int(session.get("turn", 0)) + 1
    session["updatedAt"] = now()
    course["progress"].setdefault("scores", []).append({"sceneId": args.scene, "question": args.question, "correct": correct, "at": result["at"]})
    save_course(root, course)
    print(json.dumps(result, ensure_ascii=False))


def markdown(course: dict[str, Any]) -> str:
    out = [f"# {course['title']}", "", f"**Topic:** {course['topic']}", f"**Audience:** {course['audience']}", f"**Duration:** {course['durationMinutes']} minutes", ""]
    if course.get("objectives"):
        out += ["## Learning objectives", *[f"- {x}" for x in course["objectives"]], ""]
    for n, lesson in enumerate(course.get("lessons", []), 1):
        out += [f"## {n}. {lesson.get('title', lesson['id'])}", ""]
        if lesson.get("objectives"): out += ["### Objectives", *[f"- {x}" for x in lesson["objectives"]], ""]
        for scene in lesson.get("scenes", []):
            out += [f"### {scene.get('title', scene['id'])}", "", scene.get("content", ""), ""]
        if lesson.get("exercises"):
            out.append("### Exercises")
            for ex in lesson["exercises"]:
                out += [f"- **{ex['id']} ({ex.get('type', 'short-answer')})**: {ex.get('prompt', '')}", ""]
    if course.get("sources"):
        out += ["## Sources", *[f"- {x}" for x in course["sources"]], ""]
    return "\n".join(out).rstrip() + "\n"


def html_export(course: dict[str, Any]) -> str:
    def render_scene(scene: dict[str, Any]) -> str:
        content = scene.get("content", {})
        if isinstance(content, str):
            return f"<p>{html.escape(content)}</p>"
        if content.get("type") == "interactive" and content.get("html"):
            return content["html"]
        if content.get("type") == "slide":
            bits = []
            for element in content.get("canvas", {}).get("elements", []):
                if element.get("type") == "text": bits.append(f"<p>{html.escape(str(element.get('content', '')))}</p>")
            return "".join(bits) or "<p>Slide content</p>"
        if content.get("type") == "quiz":
            bits = []
            for i, q in enumerate(content.get("questions", [])):
                answer = ",".join(str(x) for x in (q.get("answer") if isinstance(q.get("answer"), list) else [q.get("answer", "")]))
                bits.append(f"<fieldset class='quiz'><legend>{i + 1}. {html.escape(str(q.get('question', '')))}</legend><input data-answer='{html.escape(answer, quote=True)}' placeholder='Your answer'><button type='button' onclick='grade(this)'>Check</button><output></output></fieldset>")
            return "".join(bits)
        if content.get("type") == "pbl": return "<p>Project-based learning task. Continue in Hermes chat.</p>"
        return f"<p>{html.escape(str(scene.get('content', '')))}</p>"
    sections = []
    for n, lesson in enumerate(course.get("lessons", []), 1):
        scenes = "".join(f"<article><h3>{html.escape(str(s.get('title', s.get('id', 'Scene'))))}</h3>{render_scene(s)}</article>" for s in lesson.get("scenes", []))
        sections.append(f"<section id='lesson-{n}'><h2>{n}. {html.escape(str(lesson.get('title', lesson.get('id', 'Lesson'))))}</h2>{scenes}</section>")
    nav = " ".join(f"<a href='#lesson-{i}'>Lesson {i}</a>" for i in range(1, len(sections) + 1))
    sources = "".join(f"<li><a href='{html.escape(str(x.get('url', '')), quote=True)}'>{html.escape(str(x.get('title', x.get('url', ''))))}</a></li>" for x in course.get("sources", []))
    return f"<!doctype html><html><head><meta charset='utf-8'><meta name='viewport' content='width=device-width'><title>{html.escape(course['title'])}</title><style>body{{max-width:960px;margin:0 auto;padding:2rem;font:18px system-ui;line-height:1.6;color:#172033}}nav{{position:sticky;top:0;background:#fff;padding:1rem 0;border-bottom:1px solid #ddd}}nav a{{margin-right:1rem}}article{{background:#f6f8fb;padding:1rem;margin:1rem 0;border-radius:12px}}fieldset{{border:1px solid #ccd5e0;border-radius:8px;margin:1rem 0;padding:1rem}}input{{padding:.6rem;width:70%}}button{{padding:.6rem 1rem;margin-left:.5rem}}output{{display:block;margin-top:.5rem;font-weight:600}}</style></head><body><h1>{html.escape(course['title'])}</h1><p>{html.escape(course['topic'])}</p><nav>{nav}</nav>{''.join(sections)}{f'<section><h2>Sources</h2><ul>{sources}</ul></section>' if sources else ''}<script>function grade(b){{const f=b.closest('fieldset'),i=f.querySelector('input'),o=f.querySelector('output');const a=i.value.trim().toLowerCase(),e=i.dataset.answer.split(',').map(x=>x.trim().toLowerCase());o.textContent=e.includes(a)?'✓ Correct':'Review and try again';o.style.color=e.includes(a)?'green':'#a33'}}</script></body></html>\n"


def export(course: dict[str, Any], root: Path, fmt: str) -> Path:
    outdir = root / "courses" / course["id"] / "artifacts"
    outdir.mkdir(parents=True, exist_ok=True)
    if fmt == "md":
        path = outdir / f"{slug(course['title'])}.md"; path.write_text(markdown(course), encoding="utf-8")
    elif fmt == "json":
        path = outdir / "course.json"; atomic_json(path, course)
    elif fmt == "html":
        path = outdir / f"{slug(course['title'])}.html"; path.write_text(html_export(course), encoding="utf-8")
    elif fmt == "zip":
        md = outdir / f"{slug(course['title'])}.md"; md.write_text(markdown(course), encoding="utf-8")
        js = outdir / "course.json"; atomic_json(js, course)
        path = outdir / f"{slug(course['title'])}.zip"
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
            z.write(md, md.name); z.write(js, js.name)
            for artifact in course.get("artifacts", []):
                source = Path(str(artifact.get("path", ""))).expanduser()
                if source.is_file() and source.resolve() != path.resolve():
                    z.write(source, f"assets/{source.name}")
    elif fmt == "pptx":
        try:
            from pptx import Presentation
        except ImportError as error:
            raise SystemExit("optional exporter unavailable: install python-pptx") from error
        path = outdir / f"{slug(course['title'])}.pptx"
        presentation = Presentation()
        for title, content in [(course["title"], course["topic"])] + [(x.get("title", x["id"]), "\n".join(s.get("content", "") for s in x.get("scenes", []))) for x in course.get("lessons", [])]:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = content[:4000]
        presentation.save(path)
    elif fmt == "pdf":
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen.canvas import Canvas
        except ImportError as error:
            raise SystemExit("optional exporter unavailable: install reportlab") from error
        path = outdir / f"{slug(course['title'])}.pdf"
        canvas = Canvas(str(path), pagesize=letter)
        y = 760
        for line in markdown(course).splitlines():
            if y < 50:
                canvas.showPage(); y = 760
            canvas.drawString(45, y, line[:110]); y -= 14
        canvas.save()
    else:
        raise SystemExit(f"unsupported export format: {fmt}")
    course["artifacts"].append({"format": fmt, "path": str(path), "createdAt": now()})
    save_course(root, course)
    print(json.dumps({"format": fmt, "path": str(path), "bytes": path.stat().st_size}, ensure_ascii=False))
    return path


def cmd_export(args: argparse.Namespace) -> None:
    export(load_course(workspace(args), args.course), workspace(args), args.format)


def cmd_show(args: argparse.Namespace) -> None:
    print(json.dumps(load_course(workspace(args), args.course), ensure_ascii=False, indent=2))


def cmd_validate(args: argparse.Namespace) -> None:
    c = load_course(workspace(args), args.course)
    errors = []
    for key in ("id", "title", "topic", "lessons", "sessions", "progress"):
        if key not in c: errors.append(f"missing {key}")
    ids = [x.get("id") for x in c.get("lessons", [])]
    if len(ids) != len(set(ids)): errors.append("duplicate lesson ids")
    print(json.dumps({"valid": not errors, "errors": errors, "lessonCount": len(ids), "sessionCount": len(c.get('sessions', []))}, ensure_ascii=False))
    raise SystemExit(1 if errors else 0)


def parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser()
    sub = p.add_subparsers(dest="command", required=True)
    def common(x): x.add_argument("--workspace", required=True)
    x = sub.add_parser("init"); common(x); x.add_argument("--title", required=True); x.add_argument("--topic", required=True); x.add_argument("--audience", default="general adult learner"); x.add_argument("--language", default="中文"); x.add_argument("--duration", type=int, default=60); x.add_argument("--style", default="clear and practical"); x.set_defaults(func=cmd_init)
    x = sub.add_parser("lesson"); common(x); x.add_argument("--course", required=True); x.add_argument("--lesson", required=True); x.add_argument("--content", required=True); x.set_defaults(func=cmd_lesson)
    x = sub.add_parser("scene"); common(x); x.add_argument("--course", required=True); x.add_argument("--lesson", required=True); x.add_argument("--scene", required=True); x.set_defaults(func=cmd_scene)
    x = sub.add_parser("exercise"); common(x); x.add_argument("--course", required=True); x.add_argument("--lesson", required=True); x.add_argument("--exercise-id", required=True); x.add_argument("--exercise", required=True); x.set_defaults(func=cmd_exercise)
    x = sub.add_parser("session"); common(x); x.add_argument("--course", required=True); x.add_argument("--lesson"); x.set_defaults(func=cmd_session)
    x = sub.add_parser("list"); common(x); x.set_defaults(func=cmd_list)
    x = sub.add_parser("pause"); common(x); x.add_argument("--course", required=True); x.add_argument("--session", required=True); x.set_defaults(func=cmd_pause)
    x = sub.add_parser("resume"); common(x); x.add_argument("--course", required=True); x.add_argument("--session", required=True); x.set_defaults(func=cmd_resume)
    x = sub.add_parser("progress"); common(x); x.add_argument("--course", required=True); x.set_defaults(func=cmd_progress)
    x = sub.add_parser("source"); common(x); x.add_argument("--course", required=True); x.add_argument("--url", required=True); x.add_argument("--title"); x.set_defaults(func=cmd_source)
    x = sub.add_parser("artifact"); common(x); x.add_argument("--course", required=True); x.add_argument("--kind", required=True); x.add_argument("--path", required=True); x.set_defaults(func=cmd_artifact)
    x = sub.add_parser("answer"); common(x); x.add_argument("--course", required=True); x.add_argument("--session", required=True); x.add_argument("--exercise", required=True); x.add_argument("--answer", required=True); x.set_defaults(func=cmd_answer)
    x = sub.add_parser("turn"); common(x); x.add_argument("--course", required=True); x.add_argument("--session", required=True); x.add_argument("--mode", choices=["lecture", "question", "discussion", "review", "complete"], required=True); x.add_argument("--message"); x.set_defaults(func=cmd_turn)
    x = sub.add_parser("quiz-answer"); common(x); x.add_argument("--course", required=True); x.add_argument("--session", required=True); x.add_argument("--scene", required=True); x.add_argument("--question", required=True); x.add_argument("--answer", required=True); x.set_defaults(func=cmd_quiz_answer)
    x = sub.add_parser("pbl-init"); common(x); x.add_argument("--course", required=True); x.add_argument("--project", required=True); x.add_argument("--title", required=True); x.add_argument("--description", required=True); x.set_defaults(func=cmd_pbl_init)
    x = sub.add_parser("pbl-submit"); common(x); x.add_argument("--course", required=True); x.add_argument("--project", required=True); x.add_argument("--milestone", required=True); x.add_argument("--kind", choices=["text", "file", "link"], default="text"); x.add_argument("--content", required=True); x.set_defaults(func=cmd_pbl_submit)
    x = sub.add_parser("pbl-complete"); common(x); x.add_argument("--course", required=True); x.add_argument("--project", required=True); x.add_argument("--milestone", required=True); x.set_defaults(func=cmd_pbl_complete)
    x = sub.add_parser("pbl-plan"); common(x); x.add_argument("--course", required=True); x.add_argument("--project", required=True); x.add_argument("--milestones", required=True); x.set_defaults(func=cmd_pbl_plan)
    x = sub.add_parser("pbl-evaluate"); common(x); x.add_argument("--course", required=True); x.add_argument("--project", required=True); x.add_argument("--milestone", required=True); x.add_argument("--score", type=int, required=True); x.add_argument("--feedback", required=True); x.add_argument("--strength", action="append"); x.add_argument("--improvement", action="append"); x.set_defaults(func=cmd_pbl_evaluate)
    x = sub.add_parser("pbl-report"); common(x); x.add_argument("--course", required=True); x.add_argument("--project", required=True); x.set_defaults(func=cmd_pbl_report)
    x = sub.add_parser("pbl-import"); common(x); x.add_argument("--course", required=True); x.add_argument("--project-id", required=True); x.add_argument("--project", required=True); x.set_defaults(func=cmd_pbl_import)
    x = sub.add_parser("pbl-next"); common(x); x.add_argument("--course", required=True); x.add_argument("--project", required=True); x.set_defaults(func=cmd_pbl_next)
    x = sub.add_parser("pbl-task-submit"); common(x); x.add_argument("--course", required=True); x.add_argument("--project", required=True); x.add_argument("--milestone", required=True); x.add_argument("--microtask", required=True); x.add_argument("--kind", choices=["text", "file", "link"], default="text"); x.add_argument("--content", required=True); x.set_defaults(func=cmd_pbl_task_submit)
    x = sub.add_parser("pbl-task-complete"); common(x); x.add_argument("--course", required=True); x.add_argument("--project", required=True); x.add_argument("--milestone", required=True); x.add_argument("--microtask", required=True); x.set_defaults(func=cmd_pbl_task_complete)
    x = sub.add_parser("export"); common(x); x.add_argument("--course", required=True); x.add_argument("--format", choices=["md", "json", "html", "zip", "pdf", "pptx"], required=True); x.set_defaults(func=cmd_export)
    x = sub.add_parser("show"); common(x); x.add_argument("--course", required=True); x.set_defaults(func=cmd_show)
    x = sub.add_parser("validate"); common(x); x.add_argument("--course", required=True); x.set_defaults(func=cmd_validate)
    return p


if __name__ == "__main__":
    try:
        args = parser().parse_args()
        args.func(args)
    except AttributeError:
        raise SystemExit("invalid command")
